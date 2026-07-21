import io
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from openai import OpenAI

MODEL = "gpt-5.6-luna"

# ── 색상 팔레트 ───────────────────────────────────────────────────────────────
BLUE   = RGBColor(0x1F, 0x49, 0x7D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)

W, H = Inches(13.33), Inches(7.5)   # 와이드 16:9


def _add_textbox(slide, text, left, top, width, height,
                 size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def _fill_bg(slide, prs, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── 슬라이드 유형별 생성 ──────────────────────────────────────────────────────

def _slide_title(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _fill_bg(slide, prs, BLUE)
    _add_textbox(slide, title,
                 Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
                 size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _add_textbox(slide, subtitle,
                     Inches(1), Inches(4.2), Inches(11.33), Inches(0.8),
                     size=20, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)


def _slide_section(prs, title: str, bullets: list[str], score_info: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, prs, GRAY)

    hdr = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), W, Inches(1.2)
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = BLUE
    hdr.line.fill.background()

    _add_textbox(slide, title,
                 Inches(0.3), Inches(0.1), Inches(10), Inches(1),
                 size=24, bold=True, color=WHITE)
    if score_info:
        _add_textbox(slide, score_info,
                     Inches(10), Inches(0.2), Inches(3), Inches(0.8),
                     size=14, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.RIGHT)

    bullet_text = "\n".join(f"▪  {b}" for b in bullets[:7])
    _add_textbox(slide, bullet_text,
                 Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.8),
                 size=16, color=RGBColor(0x26, 0x26, 0x26))


def _slide_agenda(prs, toc: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, prs, WHITE)
    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.2))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = ACCENT
    hdr.line.fill.background()
    _add_textbox(slide, "목  차", Inches(0.3), Inches(0.1), Inches(12), Inches(1), size=26, bold=True, color=WHITE)

    items = toc.get("toc", [])
    col_size = max(1, (len(items) + 1) // 2)
    left_items = items[:col_size]
    right_items = items[col_size:]

    def fmt(sec):
        return f"{sec['order']}. {sec['title']}"

    left_text = "\n".join(fmt(s) for s in left_items)
    right_text = "\n".join(fmt(s) for s in right_items)
    _add_textbox(slide, left_text, Inches(0.5), Inches(1.5), Inches(6), Inches(5.5), size=18)
    if right_text:
        _add_textbox(slide, right_text, Inches(7), Inches(1.5), Inches(6), Inches(5.5), size=18)


# ── OpenAI: 슬라이드 불릿 + Q&A 생성 ─────────────────────────────────────────

BULLET_PROMPT = """아래 제안서 섹션 내용을 발표용 슬라이드 불릿 포인트 5~7개로 요약하세요.
각 불릿은 한 문장, 명사형 종결어미로 작성하세요.
JSON만 출력: {{"bullets": ["...", "..."]}}

섹션: {title}
내용: {content}
"""

QA_PROMPT = """당신은 공공기관 제안 평가위원입니다.
아래 제안서를 검토하고 발표 현장에서 나올 법한 예상 질문과 모범 답변을 15개 작성하세요.
JSON만 출력:
{{
  "qa": [
    {{"question": "질문 내용", "answer": "모범 답변 (3~5문장)"}}
  ]
}}

제안서 요약:
요구사항: {requirements}
평가배점: {scores}
섹션 목록: {sections}
"""


def _parse_json(raw: str) -> dict:
    cleaned = re.sub(r"```(?:json)?", "", raw).strip().strip("`").strip()
    start = cleaned.find("{")
    if start == -1:
        return {}
    end = cleaned.rfind("}") + 1
    try:
        return json.loads(cleaned[start:end])
    except json.JSONDecodeError:
        return {}


def generate_ppt(
    project_name: str,
    toc: dict,
    sections_draft: dict,
    rfp_result: dict,
    api_key: str,
) -> bytes:
    client = OpenAI(api_key=api_key)
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    _slide_title(prs, project_name, "제안발표자료")
    _slide_agenda(prs, toc)

    for sec in toc.get("toc", []):
        title   = sec["title"]
        content = sections_draft.get(title, "")

        try:
            prompt = BULLET_PROMPT.format(title=title, content=content[:3000])
            resp = client.chat.completions.create(
                model=MODEL,
                messages=[{"role": "user", "content": prompt}],
            )
            parsed = _parse_json(resp.choices[0].message.content or "")
            bullets = parsed.get("bullets", [])
        except Exception:
            bullets = [l.strip("- ").strip() for l in content.split("\n") if l.strip() and not l.startswith("#")][:6]

        if not bullets:
            bullets = [f"{sub} 관련 내용 포함" for sub in sec.get("subsections", [])]

        _slide_section(prs, title, bullets, sec.get("score_basis", ""))

    _slide_title(prs, "감사합니다", project_name)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_qa_word(
    rfp_result: dict,
    toc: dict,
    sections_draft: dict,
    api_key: str,
) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor as DRGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    client = OpenAI(api_key=api_key)
    section_titles = [s["title"] for s in toc.get("toc", [])]

    prompt = QA_PROMPT.format(
        requirements=json.dumps(rfp_result.get("requirements", [])[:15], ensure_ascii=False),
        scores=json.dumps(rfp_result.get("evaluation_scores", []), ensure_ascii=False),
        sections=json.dumps(section_titles, ensure_ascii=False),
    )

    try:
        resp = client.chat.completions.create(
            model=MODEL,
            messages=[{"role": "user", "content": prompt}],
        )
        parsed = _parse_json(resp.choices[0].message.content or "")
        qa_list = parsed.get("qa", [])
    except Exception as e:
        qa_list = [{"question": f"생성 오류: {e}", "answer": ""}]

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "맑은 고딕"
    style.font.size = Pt(10)

    t = doc.add_heading("예상 질의응답", level=1)

    for i, item in enumerate(qa_list, 1):
        q_p = doc.add_paragraph()
        q_run = q_p.add_run(f"Q{i}. {item.get('question','')}")
        q_run.bold = True
        q_run.font.color.rgb = DRGBColor(0x1F, 0x49, 0x7D)

        a_p = doc.add_paragraph()
        a_p.add_run(f"A. {item.get('answer','')}")

        doc.add_paragraph()

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()
